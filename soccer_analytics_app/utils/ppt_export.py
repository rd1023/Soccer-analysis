from pptx import Presentation
from pptx.util import Inches

def generate_powerpoint_summary(events, filename="output/presentations/match_summary.pptx"):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Match Summary - Lulu Analysis"

    for e in events[:5]:  # Limit to 5 events for preview
        content = f"{e['Timestamp']}s - {e['Event Type']} by #{e['Player']} ({e['Notes']})"
        slide.shapes.add_textbox(Inches(1), Inches(2 + 0.5 * events.index(e)), Inches(8), Inches(0.5)).text = content

    prs.save(filename)
